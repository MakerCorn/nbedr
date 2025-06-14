"""
Document processing service for chunking, text extraction, and embedding generation.
"""
import json
import logging
import re
from math import ceil
from pathlib import Path
from typing import List, Dict, Any, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed
import time

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from langchain_experimental.text_splitter import SemanticChunker
    from langchain_openai.embeddings import OpenAIEmbeddings, AzureOpenAIEmbeddings
except ImportError:
    SemanticChunker = None
    OpenAIEmbeddings = None
    AzureOpenAIEmbeddings = None

try:
    from tqdm import tqdm
except ImportError:
    def tqdm(iterable, *args, **kwargs):
        return iterable

from ..models import DocumentChunk, DocType, ChunkingStrategy, EmbeddingBatch, ProcessingResult, EmbeddingStats
from ..config import EmbeddingConfig

logger = logging.getLogger(__name__)

class DocumentService:
    """Service for document processing, chunking, and embedding generation."""
    
    def __init__(self, config: EmbeddingConfig):
        self.config = config
        self.stats = EmbeddingStats()
    
    def process_documents(self, data_path: Path) -> List[DocumentChunk]:
        """Process documents and return chunks."""
        logger.info(f"Processing documents from {data_path}")
        
        if self.config.doctype == "api":
            return self._process_api_documents(data_path)
        else:
            return self._process_regular_documents(data_path)
    
    def generate_embeddings(self, chunks: List[DocumentChunk]) -> List[DocumentChunk]:
        """Generate embeddings for document chunks."""
        logger.info(f"Generating embeddings for {len(chunks)} chunks")
        
        # TODO: Implement actual embedding generation
        # This is where you would integrate with OpenAI API or other embedding services
        embedded_chunks = []
        
        # Process in batches
        batch_size = self.config.batch_size_embeddings
        for i in range(0, len(chunks), batch_size):
            batch_chunks = chunks[i:i + batch_size]
            batch = EmbeddingBatch.create(batch_chunks, self.config.embedding_model)
            
            try:
                processed_batch = self._process_embedding_batch(batch)
                embedded_chunks.extend(processed_batch.chunks)
                batch.mark_completed()
                logger.info(f"Processed batch {i//batch_size + 1}/{ceil(len(chunks)/batch_size)}")
            except Exception as e:
                logger.error(f"Failed to process embedding batch: {e}")
                batch.mark_failed(str(e))
                # Continue with remaining batches
        
        return embedded_chunks
    
    def _process_embedding_batch(self, batch: EmbeddingBatch) -> EmbeddingBatch:
        """Process a batch of chunks for embedding generation."""
        # TODO: Implement actual embedding API calls
        # For now, return mock embeddings
        
        texts = [chunk.content for chunk in batch.chunks]
        
        # Mock embedding generation - replace with real implementation
        mock_embeddings = self._generate_mock_embeddings(texts)
        
        # Update chunks with embeddings
        for chunk, embedding in zip(batch.chunks, mock_embeddings):
            chunk.embedding = embedding
            chunk.embedding_model = batch.model
        
        return batch
    
    def _generate_mock_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Generate mock embeddings for testing purposes."""
        # TODO: Replace with actual embedding generation
        # This would call OpenAI API, Azure OpenAI, or other embedding services
        
        import random
        dimension = self.config.embedding_dimensions
        
        embeddings = []
        for text in texts:
            # Generate deterministic mock embedding based on text hash
            random.seed(hash(text) % (2**32))
            embedding = [random.uniform(-1, 1) for _ in range(dimension)]
            embeddings.append(embedding)
        
        return embeddings
    
    def store_embeddings(self, chunks: List[DocumentChunk]) -> bool:
        """Store embeddings in vector database."""
        logger.info(f"Storing {len(chunks)} embeddings in {self.config.vector_db_type} database")
        
        # TODO: Implement vector database integration
        # This would integrate with FAISS, Pinecone, Chroma, etc.
        
        try:
            # Mock storage operation
            for chunk in chunks:
                if chunk.has_embedding():
                    # Generate mock vector ID
                    chunk.vector_id = f"vec_{chunk.id}"
            
            logger.info(f"Successfully stored {len(chunks)} embeddings")
            return True
            
        except Exception as e:
            logger.error(f"Failed to store embeddings: {e}")
            return False
    
    def _process_api_documents(self, data_path: Path) -> List[DocumentChunk]:
        """Process API documentation from JSON file."""
        with open(data_path) as f:
            api_docs_json = json.load(f)
        
        required_fields = ["user_name", "api_name", "api_call", "api_version", "api_arguments", "functionality"]
        if api_docs_json and isinstance(api_docs_json[0], dict):
            for field in required_fields:
                if field not in api_docs_json[0]:
                    raise ValueError(f"API documentation missing required field: {field}")
        
        chunks = []
        for i, api_doc in enumerate(api_docs_json):
            chunk = DocumentChunk.create(
                content=str(api_doc),
                source=str(data_path),
                metadata={"type": "api", "index": i}
            )
            chunks.append(chunk)
        
        self.stats.total_chunks = len(chunks)
        return chunks
    
    def _process_regular_documents(self, data_path: Path) -> List[DocumentChunk]:
        """Process regular documents (PDF, TXT, JSON, PPTX)."""
        # Get list of files to process
        file_paths = []
        if data_path.is_dir():
            file_paths = list(data_path.rglob(f'**/*.{self.config.doctype}'))
        else:
            file_paths = [data_path]
        
        all_chunks = []
        futures = []
        
        with tqdm(total=len(file_paths), desc="Processing files", unit="file") as pbar:
            with ThreadPoolExecutor(max_workers=self.config.embed_workers) as executor:
                for file_path in file_paths:
                    future = executor.submit(self._process_single_file, file_path)
                    futures.append(future)
                    
                    if self.config.pace:
                        time.sleep(15)
                
                for future in as_completed(futures):
                    try:
                        chunks = future.result()
                        all_chunks.extend(chunks)
                        pbar.set_postfix({'total_chunks': len(all_chunks)})
                        pbar.update(1)
                    except Exception as e:
                        logger.error(f"Error processing file: {e}")
                        pbar.update(1)
        
        self.stats.total_chunks = len(all_chunks)
        return all_chunks
    
    def _process_single_file(self, file_path: Path) -> List[DocumentChunk]:
        """Process a single file and return its chunks."""
        logger.debug(f"Processing file: {file_path}")
        
        # Extract text based on document type
        text = self._extract_text(file_path)
        
        # Split into chunks
        chunk_contents = self._split_text(text)
        
        # Create DocumentChunk objects
        chunks = []
        for i, content in enumerate(chunk_contents):
            chunk = DocumentChunk.create(
                content=content,
                source=str(file_path),
                metadata={
                    "type": self.config.doctype,
                    "chunk_index": i,
                    "chunking_strategy": self.config.chunking_strategy,
                    "file_size": file_path.stat().st_size if file_path.exists() else 0
                }
            )
            chunks.append(chunk)
        
        return chunks
    
    def _extract_text(self, file_path: Path) -> str:
        """Extract text from a file based on its type."""
        if self.config.doctype == "json":
            with open(file_path, 'r') as f:
                data = json.load(f)
            return data.get("text", str(data))
        
        elif self.config.doctype == "pdf":
            text = ""
            with open(file_path, 'rb') as file:
                reader = pypdf.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text()
            return text
        
        elif self.config.doctype == "txt":
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        
        elif self.config.doctype == "pptx":
            return self._extract_text_from_pptx(file_path)
        
        else:
            raise ValueError(f"Unsupported document type: {self.config.doctype}")
    
    def _extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint file."""
        prs = Presentation(file_path)
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
                elif hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text_parts.append(cell.text)
        
        return "\n".join(text_parts)
    
    def _split_text(self, text: str) -> List[str]:
        """Split text into chunks based on the configured strategy."""
        if self.config.chunking_strategy == "semantic":
            return self._semantic_chunking(text)
        elif self.config.chunking_strategy == "fixed":
            return self._fixed_chunking(text)
        elif self.config.chunking_strategy == "sentence":
            return self._sentence_chunking(text)
        else:
            raise ValueError(f"Unknown chunking strategy: {self.config.chunking_strategy}")
    
    def _semantic_chunking(self, text: str) -> List[str]:
        """Perform semantic chunking using embeddings."""
        # TODO: Implement semantic chunking with embeddings
        # For now, fall back to fixed chunking
        logger.warning("Semantic chunking not yet implemented, using fixed chunking")
        return self._fixed_chunking(text)
    
    def _fixed_chunking(self, text: str) -> List[str]:
        """Split text into fixed-size chunks."""
        chunk_size = self.config.chunk_size
        return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    
    def _sentence_chunking(self, text: str) -> List[str]:
        """Split text by sentences, respecting chunk size limits."""
        sentences = re.split(r'(?<=[.!?]) +', text)
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            if len(current_chunk) + len(sentence) > self.config.chunk_size:
                if current_chunk:
                    chunks.append(current_chunk)
                current_chunk = sentence
            else:
                current_chunk += (" " if current_chunk else "") + sentence
        
        if current_chunk:
            chunks.append(current_chunk)
        
        return chunks
    
    def get_stats(self) -> EmbeddingStats:
        """Get processing statistics."""
        return self.stats